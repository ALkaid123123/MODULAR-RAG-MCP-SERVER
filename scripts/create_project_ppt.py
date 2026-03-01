#!/usr/bin/env python3
"""
Generate project presentation (PPTX) for Modular RAG MCP Server.
Run: pip install python-pptx && python scripts/create_project_ppt.py
Output: docs/Modular-RAG-MCP-Server.pptx (or ./Modular-RAG-MCP-Server.pptx)
"""

from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError as e:
    raise SystemExit(f"Install python-pptx: pip install python-pptx\n{e}") from e

# Theme: Teal Trust / Ocean — fits RAG/AI/tech
TITLE_BG = RGBColor(0x02, 0x80, 0x90)   # teal
SECONDARY = RGBColor(0x00, 0xA8, 0x96)  # seafoam
ACCENT = RGBColor(0x02, 0xC3, 0x9A)     # mint
DARK = RGBColor(0x21, 0x29, 0x5C)       # midnight
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x55, 0x55, 0x55)


def set_title_slide(prs, title: str, subtitle: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    shp = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = TITLE_BG
    shp.line.fill.background()
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), prs.slide_width - Inches(1), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), prs.slide_width - Inches(1), Inches(0.8))
        sp = sub_box.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(22)
        sp.font.color.rgb = RGBColor(0xE0, 0xFF, 0xFF)
        sp.alignment = PP_ALIGN.CENTER
    return slide


def add_section_slide(prs, title: str, bullets: list[str], notes: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Title bar
    bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK
    bar.line.fill.background()
    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), prs.slide_width - Inches(1), Inches(0.7))
    tp = tbox.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(28)
    tp.font.bold = True
    tp.font.color.rgb = WHITE
    # Body
    body = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), prs.slide_width - Inches(1), Inches(5.5))
    tf = body.text_frame
    tf.word_wrap = True
    for i, line in enumerate(bullets):
        p = tf.add_paragraph() if i else tf.paragraphs[0]
        p.text = line
        p.font.size = Pt(18)
        p.font.color.rgb = GRAY
        p.space_after = Pt(10)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def add_two_column_slide(prs, title: str, left_bullets: list[str], right_bullets: list[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK
    bar.line.fill.background()
    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), prs.slide_width - Inches(1), Inches(0.7))
    tbox.text_frame.paragraphs[0].text = title
    tbox.text_frame.paragraphs[0].font.size = Pt(28)
    tbox.text_frame.paragraphs[0].font.bold = True
    tbox.text_frame.paragraphs[0].font.color.rgb = WHITE
    w = (prs.slide_width - Inches(1.5)) / 2
    left = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), w, Inches(5))
    right = slide.shapes.add_textbox(Inches(0.5) + w + Inches(0.5), Inches(1.4), w, Inches(5))
    for i, line in enumerate(left_bullets):
        p = left.text_frame.add_paragraph() if i else left.text_frame.paragraphs[0]
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = GRAY
    for i, line in enumerate(right_bullets):
        p = right.text_frame.add_paragraph() if i else right.text_frame.paragraphs[0]
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = GRAY
    return slide


def main():
    out_dir = Path(__file__).resolve().parents[1]
    docs = out_dir / "docs"
    docs.mkdir(exist_ok=True)
    out_path = docs / "Modular-RAG-MCP-Server.pptx"

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 1. Title
    set_title_slide(
        prs,
        "Modular RAG MCP Server",
        "可插拔、可观测的模块化 RAG 服务 · 通过 MCP 对接 Copilot / Claude"
    )

    # 2. 项目概述
    add_section_slide(prs, "项目概述", [
        "基于 RAG（检索增强生成）与 MCP（Model Context Protocol）的智能问答与知识检索框架",
        "目标：可扩展、高可观测、易迭代",
        "定位：自学与教学同步（Learning by Teaching）",
        "面向 RAG 技术学习与面试求职的实战平台",
    ], "设计理念：教是最好的学")

    # 3. 核心价值
    add_section_slide(prs, "核心价值", [
        "实战驱动：架构即 RAG 面试题活体答案（分层检索、Hybrid Search、Rerank、评测）",
        "开箱即用 + 深度扩展：MCP 标准接口对接 Copilot/Claude，内部完全模块化可替换",
        "配套资源：技术文档、带注释源码、视频讲解、知识点清单与面试题",
        "可观测与可量化：全链路追踪、Streamlit 管理面板、Ragas/Custom 评估闭环",
    ])

    # 4. 系统架构概览
    add_section_slide(prs, "系统架构概览", [
        "Ingestion Pipeline：PDF → Markdown → Chunk → Transform → Embedding → Upsert（含多模态图片描述）",
        "Hybrid Search：Dense（向量）+ Sparse（BM25）+ RRF Fusion + 可选 Rerank",
        "MCP Server：Stdio 传输，暴露 query_knowledge_hub、list_collections、get_document_summary",
        "Dashboard：Streamlit 六页面（总览 / 数据浏览 / Ingestion / 追踪 / 评估）",
        "Evaluation：Ragas + Custom，支持 golden test set 回归",
    ])

    # 5. RAG 策略亮点
    add_section_slide(prs, "RAG 策略亮点", [
        "分块：智能分块 + 上下文增强（元数据、Image Caption）",
        "粗排：混合检索（BM25 + Dense Embedding）+ RRF 融合，平衡查全与查准",
        "精排：Cross-Encoder 或 LLM Rerank，粗排→精排两段式，兼顾速度与 Top 精准度",
    ])

    # 6. 可插拔架构
    add_two_column_slide(prs, "可插拔架构", [
        "LLM：Azure / OpenAI / Ollama / DeepSeek 等，配置一键切换",
        "Embedding & Rerank：云端与本地模型统一接口",
        "Loader / Splitter / Transform：PDF、Markdown、语义切分、Image Caption 可替换",
        "检索策略：纯向量 / 纯关键词 / 混合可配置",
        "向量库：Chroma、Qdrant、Milvus 等可换",
        "评估：Ragas、DeepEval 等可挂载",
    ], [
        "零代码修改即可 A/B 测试、成本优化、隐私迁移",
        "接口隔离 + 配置驱动 + 工厂模式 + 优雅降级",
    ])

    # 7. MCP 生态集成
    add_section_slide(prs, "MCP 生态集成", [
        "Server 暴露标准 tools/resources，Copilot、Claude Desktop 等 MCP Client 直接连接",
        "Stdio 本地通信：零端口、无鉴权、数据不出本机，适合私有知识库",
        "零前端开发：复用 VS Code / 编辑器与现有 AI 助手",
        "一次开发，处处可用：任何支持 MCP 的 Agent 均可接入",
    ])

    # 8. 多模态与可观测
    add_section_slide(prs, "多模态与可观测", [
        "多模态：Image-to-Text 策略，Vision LLM 生成图像描述并写入 Chunk，统一向量空间检索",
        "全链路白盒：Ingestion 与 Query 各阶段可追踪、可可视化",
        "Dashboard：系统总览、数据浏览、Ingestion 管理、Query/Ingestion 追踪、评估面板",
        "评估闭环：Hit Rate、MRR、Faithfulness 等指标，数据驱动迭代",
    ])

    # 9. 技术栈
    add_section_slide(prs, "技术栈（选型要点）", [
        "配置：YAML（settings.yaml），统一配置驱动",
        "PDF→Markdown：MarkItDown；切分：LangChain RecursiveCharacterTextSplitter",
        "向量库：Chroma；Embedding：OpenAI / Azure / Ollama 等",
        "MCP：Python 官方 mcp SDK，Stdio Transport",
        "持久化：SQLite（ingestion_history、image_index 等），本地优先",
    ])

    # 10. 项目状态与许可
    set_title_slide(
        prs,
        "项目状态与许可",
        "开发中 (WIP) · 预计 2026 年 3 月完成 · MIT License"
    )

    prs.save(str(out_path))
    print(f"Saved: {out_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
